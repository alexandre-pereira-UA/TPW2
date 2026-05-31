import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def gerar_apresentacao():
    print("A iniciar a geração do PowerPoint...")
    prs = Presentation()
    
    # Configura o tamanho dos slides para Widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Layout em branco para total controle dos shapes
    blank_layout = prs.slide_layouts[6]
    
    # -------------------------------------------------------------
    # SLIDE 1: Capa de Apresentação (Estilo Escuro / Gold Cinema)
    # -------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    
    # Adiciona fundo preto/indigo profundo
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height) # MSO_SHAPE.RECTANGLE = 1
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(9, 9, 13)
    bg.line.fill.background()
    
    # Caixa de título da Capa
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(4.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "MOVIEZ CATALOG"
    p1.font.name = 'Inter'
    p1.font.size = Pt(56)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(245, 197, 24) # Amarelo Ouro
    
    p2 = tf.add_paragraph()
    p2.text = "Plataforma Web Desacoplada | Trabalho Prático 2"
    p2.font.name = 'Inter'
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(255, 255, 255)
    p2.space_before = Pt(10)
    
    p3 = tf.add_paragraph()
    p3.text = "Grupo de Trabalho:\n• Guilherme Escórcio — Nº 118648\n• Alexandre Pereira — Nº 119871\n• Alexandre Silva — Nº 119583"
    p3.font.name = 'Inter'
    p3.font.size = Pt(15)
    p3.font.color.rgb = RGBColor(210, 210, 225)
    p3.space_before = Pt(25)
    
    p4 = tf.add_paragraph()
    p4.text = "Tecnologias e Programação Web (TPW) | Engenharia Informática, Universidade de Aveiro\nAno Letivo: 2025/2026"
    p4.font.name = 'Inter'
    p4.font.size = Pt(12)
    p4.font.color.rgb = RGBColor(140, 140, 160)
    p4.space_before = Pt(30)
    
    # -------------------------------------------------------------
    # HELPER: Adicionar Slides de Conteúdo com Estilo Consistente e Split-Screen
    # -------------------------------------------------------------
    def add_content_slide(title_text, items, image_name=None):
        s = prs.slides.add_slide(blank_layout)
        
        # Fundo cinemático
        b = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        b.fill.solid()
        b.fill.fore_color.rgb = RGBColor(18, 18, 28) # Deep indigo gray
        b.line.fill.background()
        
        # Título do Slide
        t_box = s.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(1.0))
        tf_title = t_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = 'Inter'
        p_title.font.size = Pt(32)
        p_title.font.bold = True
        p_title.font.color.rgb = RGBColor(245, 197, 24)
        
        # Definição das dimensões da caixa de texto conforme exista ou não imagem
        if image_name:
            text_width = Inches(6.2)
            img_left = Inches(7.4)
            img_top = Inches(1.8)
            img_width = Inches(5.1)
            img_height = Inches(4.2)
        else:
            text_width = Inches(11.733)
            
        # Caixa de Texto Principal (Bullets)
        b_box = s.shapes.add_textbox(Inches(0.8), Inches(1.7), text_width, Inches(5.0))
        tf_body = b_box.text_frame
        tf_body.word_wrap = True
        
        for idx, item in enumerate(items):
            if idx == 0:
                p_item = tf_body.paragraphs[0]
            else:
                p_item = tf_body.add_paragraph()
            
            p_item.font.name = 'Inter'
            p_item.space_after = Pt(8)
            
            if item.startswith("    -"):
                # Sub-item (Nível 2)
                p_item.text = item.strip().lstrip("-").strip()
                p_item.level = 2
                p_item.font.size = Pt(14)
                p_item.font.color.rgb = RGBColor(160, 160, 180)
            elif item.startswith("  -") or item.startswith("  *"):
                # Item Secundário (Nível 1)
                p_item.text = item.strip().lstrip("-").lstrip("*").strip()
                p_item.level = 1
                p_item.font.size = Pt(16)
                p_item.font.color.rgb = RGBColor(210, 210, 225)
            else:
                # Tópico Principal (Nível 0)
                p_item.text = item.strip().lstrip("-").lstrip("*").strip()
                p_item.level = 0
                p_item.font.size = Pt(18)
                p_item.font.bold = True
                p_item.font.color.rgb = RGBColor(255, 255, 255)
                
        # Adiciona a Imagem se ela existir no disco
        if image_name:
            caminho_imagem = os.path.join("apresentacao_imagens", image_name)
            if os.path.exists(caminho_imagem):
                s.shapes.add_picture(caminho_imagem, img_left, img_top, img_width, img_height)
            else:
                print(f"Aviso: Imagem '{caminho_imagem}' não encontrada. Slide criado apenas com texto.")

    # -------------------------------------------------------------
    # SLIDE 2: Introdução e Arquitetura
    # -------------------------------------------------------------
    add_content_slide(
        "1. Introdução e Arquitetura do Sistema",
        [
            "Evolução para N-Tier Desacoplada (TP2)",
            "  - Transição de templates no servidor para duas camadas autónomas.",
            "  - Comunicação baseada exclusivamente em protocolo HTTP trocando JSON.",
            "Separação Estrita de Responsabilidades",
            "  - Servidor: API REST pura desenvolvida em Django REST Framework.",
            "  - Cliente: Single Page Application (SPA) reativa construída com Angular.",
            "Mais-Valias de Engenharia",
            "  - Modularidade do código, escalabilidade e altíssima velocidade de renderização."
        ],
        "slide2.png"
    )
    
    # -------------------------------------------------------------
    # SLIDE 3: Back-end: Django REST Framework (DRF)
    # -------------------------------------------------------------
    add_content_slide(
        "2. Back-end: Robustez e RESTful APIs",
        [
            "Modelo de Dados Otimizado (SQLite3)",
            "  - CRUD completo para Filmes, Atores, Realizadores e Géneros.",
            "  - Relações N-para-N para Criticas/Avaliações e listas personalizadas.",
            "Camada de Serialização Avançada",
            "  - Nested Serializers de Leitura: Traz dados aninhados para poupar pedidos HTTP.",
            "  - Serializers de Escrita: IDs de chave primária diretos para inserções velozes.",
            "Endpoints REST Otimizados (/ws/)",
            "  - Controlo transacional nativo com filtragem robusta e resposta em JSON puro."
        ],
        "slide3.png"
    )
    
    # -------------------------------------------------------------
    # SLIDE 4: Front-end: Angular SPA
    # -------------------------------------------------------------
    add_content_slide(
        "3. Front-end: SPA Moderna e Reativa",
        [
            "Arquitetura Reativa com Angular (v21+)",
            "  - Componentes Standalone livres de módulos e injeção de dependências nativa.",
            "  - Serviços dedicados para isolar toda a lógica de comunicação externa.",
            "Scroll Infinito Progressivo (Performance)",
            "  - Diretiva HostListener para carregar e renderizar os cards em lotes de 6 em 6.",
            "Funcionalidades Visuais Premium",
            "  - Modo Grelha Simétrica ou Modo Carrossel Automático com blur dinâmico.",
            "  - Tratamento inteligente e centrado de estado vazio de pesquisa (sem resultados)."
        ],
        "slide4.png"
    )
    
    # -------------------------------------------------------------
    # SLIDE 5: Alojamento e Integração TMDB
    # -------------------------------------------------------------
    add_content_slide(
        "4. Alojamento em Produção e TMDB",
        [
            "Alojamento Distribuído em Produção Cloud",
            "  - API Back-end: Publicada e a correr ativamente no PythonAnywhere.",
            "  - SPA Front-end: Publicada e otimizada na plataforma cloud Vercel.",
            "Integração de Valorização (TMDB API)",
            "  - Rota administrativa dedicada para consumir a API externa do The Movie Database.",
            "Importação e Normalização Dinâmica",
            "  - Descarrega posters em alta qualidade, sinopses e detalhes sob demanda.",
            "  - Cria registos locais sem duplicar realizadores ou atores existentes."
        ],
        "slide5.png"
    )
    
    # -------------------------------------------------------------
    # SLIDE 6: Segurança, Perfis e Acesso
    # -------------------------------------------------------------
    add_content_slide(
        "5. Segurança, Autenticação e Perfis",
        [
            "Autenticação por Token Segura",
            "  - TokenAuthentication integrado nas rotas protegidas e guardado em localStorage.",
            "Controlo de Acessos por Perfis de Utilizador",
            "  - Administrador (Superuser): Gestão total de estatísticas, grupos e importador TMDB.",
            "  - Moderador (Membro de 'coment'): Exclusão rápida de comentários impróprios.",
            "  - Utilizador Registado: Gere listas de Favoritos/Guardados e deixa notas por estrelas.",
            "  - Visitante Anónimo: Apenas leitura segura do catálogo de filmes."
        ],
        "slide6.png"
    )
    
    # -------------------------------------------------------------
    # SLIDE 7: Conclusão
    # -------------------------------------------------------------
    add_content_slide(
        "6. Conclusão",
        [
            "Conformidade Absoluta com o Guião Prático",
            "  - Todos os requisitos mínimos e tópicos de exploração concluídos com excelência.",
            "Plataforma Pronta para Demonstração Imediata",
            "  - Base de dados SQLite3 pré-populada com 50 filmes de grande sucesso.",
            "  - 5 contas prontas com perfis de utilizador distintos pré-configurados.",
            "  - Relatório técnico simplificado RELATORIO.md com capa oficial incluído.",
            "Impacto do Desacoplamento N-Tier",
            "  - Base sólida, modular, segura e preparada para futuras evoluções de mercado."
        ]
    )
    
    # Grava o ficheiro de apresentação no disco
    nome_arquivo = "APRESENTACAO_TP2.pptx"
    prs.save(nome_arquivo)
    print(f"Sucesso! Apresentação PowerPoint gerada em '{nome_arquivo}'.")

if __name__ == "__main__":
    gerar_apresentacao()
